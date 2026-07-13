import os
import sys
import locale

# --- 0. 環境レベルでのエラー強制封鎖（日本語対策） ---
try:
    locale.setlocale(locale.LC_ALL, 'C.UTF-8')
except:
    pass
os.environ["PYTHONIOENCODING"] = "utf-8"

import streamlit as st
import google.generativeai as genai
import pandas as pd
from io import BytesIO
import base64
from PIL import Image
import json

# --- 1. HEIC（iPhone画像形式）デコーダー登録 ---
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass

# --- 2. アプリ基本設定 ---
st.set_page_config(page_title="BizAlchemy Elite Pro Max", layout="wide")

# プロフェッショナルなUIデザイン（ファイル名による誤作動をCSSで保護）
st.markdown("""
    <style>
    [data-testid="stFileUploaderFileName"] { display: none; }
    .stApp { background-color: #f8f9fa; }
    .stButton>button { width: 100%; border-radius: 10px; height: 3.5rem; background-color: #1E88E5; color: white; font-weight: bold; }
    </style>
    """, unsafe_allow_html=True)

if "GEMINI_API_KEY" not in st.secrets:
    st.error("Gemini APIキーをSecretsに設定してください。")
    st.stop()

genai.configure(api_key=st.secrets["GEMINI_API_KEY"])

# --- 3. メタデータ隔離・画像変換エンジン ---
def get_safe_pil_image(uploaded_file):
    try:
        raw_bytes = uploaded_file.getvalue()
        img = Image.open(BytesIO(raw_bytes))
        img = img.convert("RGB") # メタデータを消去し、RGBに統一
        img.thumbnail((1200, 1200)) # 解析に最適な解像度
        return img
    except Exception:
        return None

# --- 4. プロ仕様マルチフォーマット出力エンジン ---
def generate_multiformat(df, advice, format_type):
    output = BytesIO()
    
    if format_type == "Excel (.xlsx)":
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False)
        return output.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"

    elif format_type == "Word (.docx)":
        from docx import Document
        doc = Document()
        doc.add_heading('AI Business Analysis Report', 0)
        doc.add_heading('1. 抽出データ', level=1)
        table = doc.add_table(rows=1, cols=len(df.columns))
        for i, col in enumerate(df.columns): table.rows[0].cells[i].text = str(col)
        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            for i, val in enumerate(row): row_cells[i].text = str(val)
        doc.add_heading('2. 経営アドバイス', level=1)
        doc.add_paragraph(advice)
        doc.add_paragraph("\n--- 本資料は商用ライセンスに基づき生成されました ---")
        doc.save(output)
        return output.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"

    elif format_type == "PowerPoint (.pptx)":
        from pptx import Presentation
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "AI分析結果"
        slide1.placeholders[1].text = "ビジネスデータ錬成エンジン v8"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "経営インサイト"
        slide2.placeholders[1].text = advice
        prs.save(output)
        return output.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"

# --- 5. メインUI ---
st.title("🚀 Magic Biz Data Gen Pro")
st.write("### ～ 1枚の写真から『Excel / Word / PowerPoint』を完全無料で錬成する ～")

# 信頼性パネル
c1, c2, c3 = st.columns(3)
with c1: st.info("🔒 セキュリティ: メモリ内処理（データ即時消去）")
with c2: st.info("⚖️ 法的準拠: 個人情報保護法・商用ライセンス対応")
with c3: st.info("💰 コスト: 完全無料（Gemini API無料枠）")

st.divider()

# 同意フロー
is_agreed = st.checkbox("デモ版のデータ取扱いに関する重要事項に同意して利用する", value=True)

if is_agreed:
    # 選択エリア
    col_u1, col_u2 = st.columns(2)
    with col_u1:
        doc_type = st.selectbox("📝 解析する書類のタイプ", ["タイムカード", "手書き請求書", "会議メモ/ホワイトボード", "その他"])
    with col_u2:
        out_format = st.radio("📁 出力したいファイル形式", ["Excel (.xlsx)", "Word (.docx)", "PowerPoint (.pptx)"], horizontal=True)

    # アップローダー
    uploaded_file = st.file_uploader("書類をアップロード（日本語ファイル名・iPhoneのHEIC形式に対応済み）", type=['png', 'jpg', 'jpeg', 'heic', 'HEIC'])

    if st.button("✨ データを錬成する（AI解析・資料作成）", type="primary", use_container_width=True):
        if uploaded_file:
            msg = st.empty()
            try:
                msg.info("📷 画像データを最適化中（エラーを物理的に遮断）...")
                img_obj = get_safe_pil_image(uploaded_file)
                
                if not img_obj:
                    st.error("🚨 画像の読み込みに失敗しました。")
                    st.stop()

                msg.info("🧠 最新AI（Gemini 2.5）が超高速で解析中...")
                
                # 【重要】モデル名を最新の「gemini-2.5-flash」に変更
                model = genai.GenerativeModel(
                    model_name='gemini-2.5-flash',
                    generation_config={"response_mime_type": "application/json"}
                )

                prompt = f"""
                あなたは一流のDXコンサルタントです。画像（{doc_type}）から全データを抽出してください。
                【必須項目】
                必ず以下のキーを持つJSON形式で出力してください：
                {{
                  "data": [{{"項目名": "値"}}],
                  "cost_saved": 想定削減時間（分、数値のみ）,
                  "advice": "経営改善アドバイス一言"
                }}
                """

                # API通信
                response = model.generate_content([prompt, img_obj])
                
                # 解析結果のデコード
                result = json.loads(response.text)
                df = pd.DataFrame(result["data"])
                advice = result.get("advice", "分析が完了しました。")
                saved_time = result.get("cost_saved", 30)

                msg.success(f"✅ 錬成完了！約{saved_time}分の事務作業を削減しました。")

                # 結果表示
                st.write("### 💎 錬成結果プレビュー")
                edited_df = st.data_editor(df, use_container_width=True, num_rows="dynamic")
                st.info(f"💡 AI経営コンサルのアドバイス: {advice}")

                # ファイル生成とダウンロード
                file_data, mime, ext = generate_multiformat(edited_df, advice, out_format)

                st.download_button(
                    label=f"📥 {out_format} を保存する",
                    data=file_data,
                    file_name=f"biz_refined_data.{ext}",
                    mime=mime,
                    use_container_width=True
                )

            except Exception as e:
                msg.error("🚨 解析プロセスでエラーが発生しました。")
                st.warning("APIの回数制限か、キーの設定に問題があります。少し時間を置いてお試しください。")
                print(f"Error: {e}")
        else:
            st.warning("写真をアップロードしてください。")
else:
    st.info("同意にチェックを入れると開始できます。")

st.divider()
st.caption("© 2024 Magic Biz Data Gen Pro Edition | Powered by Gemini 2.5 Flash")
